from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
SCREENSHOT_DIR = DOCS_DIR / "screenshots"
OUTPUT_PPTX = DOCS_DIR / "Group6_Presentation.pptx"
OUTPUT_PPTX_EN = DOCS_DIR / "Group6_Presentation_English.pptx"

BLUE = RGBColor(29, 78, 216)
DARK = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
LIGHT_BG = RGBColor(248, 250, 252)
GOLD = RGBColor(245, 158, 11)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(203, 213, 225)
DEFAULT_FONT = "Calibri"


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
      shape.line.color.rgb = line
    else:
      shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    size=20,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    font_name=DEFAULT_FONT,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Pt(2)
    frame.margin_right = Pt(2)
    frame.margin_top = Pt(1)
    frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.name = font_name
    font.color.rgb = color
    return textbox


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    font_size=21,
    color=DARK,
    font_name=DEFAULT_FONT,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(2)
    frame.margin_right = Pt(2)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.name = font_name
        paragraph.space_after = Pt(8)
    return textbox


def add_header(slide, title, subtitle="", slide_no=None):
    add_rect(slide, 0, 0, Inches(13.333), Inches(0.7), BLUE)
    add_textbox(slide, Inches(0.45), Inches(0.12), Inches(9.6), Inches(0.35), title, size=27, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.48), Inches(0.78), Inches(10.2), Inches(0.3), subtitle, size=11, color=SLATE)
    if slide_no is not None:
        add_textbox(
            slide,
            Inches(12.35),
            Inches(0.15),
            Inches(0.55),
            Inches(0.3),
            str(slide_no),
            size=16,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.RIGHT,
        )


def add_picture_contain(slide, path: Path, left, top, width, height):
    img = Image.open(path)
    img_ratio = img.width / img.height
    box_ratio = width / height

    if img_ratio > box_ratio:
        render_width = width
        render_height = width / img_ratio
        render_left = left
        render_top = top + (height - render_height) / 2
    else:
        render_height = height
        render_width = height * img_ratio
        render_top = top
        render_left = left + (width - render_width) / 2

    slide.shapes.add_picture(str(path), render_left, render_top, width=render_width, height=render_height)


def add_caption(slide, left, top, width, text):
    add_textbox(slide, left, top, width, Inches(0.36), text, size=10.5, color=SLATE, align=PP_ALIGN.CENTER)


def add_card(slide, left, top, width, height, title, body, accent=BLUE):
    add_rect(slide, left, top, width, height, WHITE, line=BORDER, radius=True)
    add_rect(slide, left, top, Inches(0.08), height, accent)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.3), Inches(0.32), title, size=16, bold=True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.48), width - Inches(0.28), height - Inches(0.58), body, size=11.5, color=SLATE)


def add_user_card(slide, left, top, width, height, title, bullets, accent):
    add_rect(slide, left, top, width, height, WHITE, line=BORDER, radius=True)
    add_rect(slide, left, top, width, Inches(0.18), accent)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.24), width - Inches(0.24), Inches(0.5), title, size=14, bold=True)
    add_bullets(slide, left + Inches(0.16), top + Inches(0.9), width - Inches(0.24), height - Inches(1.0), bullets, font_size=11, color=SLATE)


def add_arch_box(slide, left, top, width, height, text, fill):
    box = add_rect(slide, left, top, width, height, fill, line=BLUE, radius=True)
    add_textbox(slide, left, top + Inches(0.18), width, Inches(0.35), text, size=16, bold=True, align=PP_ALIGN.CENTER)
    return box


def add_arrow(slide, x1, y1, x2, y2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = BLUE
    connector.line.width = Pt(2.2)
    connector.line.end_arrowhead = True
    return connector


def add_table(slide, left, top, width, height, rows, cols):
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.first_row = True
    for col_idx in range(cols):
        table.columns[col_idx].width = int(width / cols)
    return table


def new_slide(prs: Presentation, layout):
    return prs.slides.add_slide(layout)


def save_with_fallback(prs: Presentation, path: Path) -> Path:
    candidates = [path]
    candidates.extend(path.with_stem(f"{path.stem}_updated_{idx}") for idx in range(1, 6))
    for candidate in candidates:
        try:
            prs.save(str(candidate))
            return candidate
        except PermissionError:
            continue
    raise PermissionError(f"Could not save presentation to {path} or fallback filenames.")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    screenshots = {
        "home": SCREENSHOT_DIR / "home_page.png",
        "search": SCREENSHOT_DIR / "search_page.png",
        "detail": SCREENSHOT_DIR / "listing_detail.png",
        "form": SCREENSHOT_DIR / "listing_form_pricing.png",
        "messages": SCREENSHOT_DIR / "messages_page.png",
        "meetup": SCREENSHOT_DIR / "meetup_page.png",
        "admin": SCREENSHOT_DIR / "admin_dashboard.png",
    }

    # Slide 1
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.18), BLUE)
    add_textbox(slide, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.35),
                "EduShare: A Smart Web-Based Marketplace for University Students to Exchange Textbooks and Learning Materials",
                size=26, bold=True)
    add_textbox(slide, Inches(0.92), Inches(2.45), Inches(7.2), Inches(0.5),
                "[CAP126] New Programming Language", size=18, bold=True, color=BLUE)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(3.15),
        Inches(6.6),
        Inches(2.0),
        [
            "Lecturer: Hannah Vu",
            "Group: Group 6",
            "Student: Mai Tan Phat",
            "Student ID: 2280602300",
            "Class: 22DTHQA2",
        ],
        font_size=18,
        color=SLATE,
    )
    add_picture_contain(slide, screenshots["home"], Inches(8.1), Inches(1.55), Inches(4.1), Inches(3.0))
    add_rect(slide, Inches(0.95), Inches(5.85), Inches(1.95), Inches(0.42), BLUE, radius=True)
    add_textbox(slide, Inches(1.1), Inches(5.94), Inches(1.65), Inches(0.22), "MERN Project", size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.1), Inches(5.85), Inches(2.45), Inches(0.42), GOLD, radius=True)
    add_textbox(slide, Inches(3.25), Inches(5.94), Inches(2.15), Inches(0.22), "Recommendation", size=10.25, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5.75), Inches(5.85), Inches(1.35), Inches(0.42), BLUE, radius=True)
    add_textbox(slide, Inches(5.9), Inches(5.94), Inches(1.05), Inches(0.22), "Safety", size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.95), Inches(6.45), Inches(11.0), Inches(0.25),
                "Ho Chi Minh City, March 2026", size=12, color=SLATE)

    # Slide 2
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "Problem and Motivation", "Why this project matters on campus", 2)
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.3),
        Inches(6.0),
        Inches(4.8),
        [
            "Students often need affordable second-hand textbooks and study materials for one semester only.",
            "Most exchanges happen in scattered Facebook posts or chat groups, which are hard to search and manage.",
            "Buyers and sellers need a safer way to communicate and arrange meetup locations on campus.",
            "A student marketplace should support course codes, trust signals, and guided meetup planning.",
        ],
        font_size=22,
    )
    add_card(
        slide,
        Inches(7.15),
        Inches(1.45),
        Inches(5.4),
        Inches(3.9),
        "Project Motivation",
        "EduShare was designed as a practical campus marketplace. Instead of building a simple CRUD website, the project focuses on recommendation, trust, moderation, and safer transactions for real student use cases.",
        accent=GOLD,
    )
    add_textbox(slide, Inches(7.3), Inches(5.65), Inches(4.9), Inches(0.6),
                "Goal: build a realistic student-focused web application, not just a listing board.", size=16, bold=True, color=BLUE)

    # Slide 3
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "Objectives and Target Users", "What the system is expected to achieve", 3)
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.25),
        Inches(5.6),
        Inches(4.9),
        [
            "Build a web-based marketplace for textbooks, notes, lab kits, and learning materials.",
            "Support registration, login, profile management, and protected student actions.",
            "Allow users to create, search, reserve, and manage listings.",
            "Provide messaging, meetup planning, reviews, and notifications in one workflow.",
            "Add recommendation, price suggestion, and moderation features to improve usability and safety.",
        ],
        font_size=18.5,
    )
    add_user_card(slide, Inches(6.3), Inches(1.35), Inches(1.8), Inches(4.65), "Guest", ["Browse listings", "Search marketplace", "View item details"], BLUE)
    add_user_card(slide, Inches(8.25), Inches(1.35), Inches(2.55), Inches(4.65), "Registered User", ["Create listings", "Reserve items", "Message sellers", "Plan meetups", "Write reviews"], GOLD)
    add_user_card(slide, Inches(10.95), Inches(1.35), Inches(1.8), Inches(4.65), "Admin", ["View overview", "Moderate users", "Handle reports", "Manage reviews"], BLUE)

    # Slide 4
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "System Overview", "Main workflow from browsing to transaction completion", 4)
    add_arch_box(slide, Inches(0.85), Inches(1.65), Inches(2.1), Inches(0.75), "Browse Listings", WHITE)
    add_arch_box(slide, Inches(3.25), Inches(1.65), Inches(2.2), Inches(0.75), "Reserve and Chat", WHITE)
    add_arch_box(slide, Inches(5.8), Inches(1.65), Inches(2.2), Inches(0.75), "Plan Meetup", WHITE)
    add_arch_box(slide, Inches(8.35), Inches(1.65), Inches(2.25), Inches(0.75), "Complete Transaction", WHITE)
    add_arch_box(slide, Inches(10.95), Inches(1.65), Inches(1.5), Inches(0.75), "Review", WHITE)
    add_arrow(slide, Inches(2.95), Inches(2.02), Inches(3.23), Inches(2.02))
    add_arrow(slide, Inches(5.48), Inches(2.02), Inches(5.78), Inches(2.02))
    add_arrow(slide, Inches(8.03), Inches(2.02), Inches(8.33), Inches(2.02))
    add_arrow(slide, Inches(10.63), Inches(2.02), Inches(10.93), Inches(2.02))
    add_picture_contain(slide, screenshots["home"], Inches(0.8), Inches(3.0), Inches(5.8), Inches(3.5))
    add_card(slide, Inches(6.95), Inches(3.1), Inches(5.3), Inches(3.3), "Key Idea",
             "EduShare combines product discovery, trust signals, reservation control, campus meetup guidance, and moderation in one student-focused workflow. This makes the system closer to a real marketplace product than a basic listing application.", accent=BLUE)

    # Slide 5
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "Core Features", "Main functions delivered in the final system", 5)
    feature_cards = [
        ("Authentication", "Register, login, forgot password, reset password, and protected routes."),
        ("Listing CRUD", "Create, edit, browse, reserve, release, and manage study material listings."),
        ("Search and Filter", "Search by keyword, category, condition, price, location, and sorting."),
        ("Messaging", "Listing-based chat between buyers, sellers, and admin support."),
        ("Meetup Planner", "Campus-safe location suggestions with date, time, and confirmation flow."),
        ("Reviews and Notifications", "Post-transaction reviews, wishlist tracking, and alert center."),
    ]
    x_positions = [Inches(0.7), Inches(4.45), Inches(8.2)]
    y_positions = [Inches(1.45), Inches(4.05)]
    for idx, (title, body) in enumerate(feature_cards):
        x = x_positions[idx % 3]
        y = y_positions[idx // 3]
        add_card(slide, x, y, Inches(3.3), Inches(2.15), title, body, accent=BLUE if idx % 2 == 0 else GOLD)

    # Slide 6
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "Smart Features", "How the project goes beyond basic CRUD", 6)
    add_card(slide, Inches(0.7), Inches(1.45), Inches(3.9), Inches(4.6), "Recommendation System",
             "Ranks listings using keyword overlap, course code match, browsing similarity, price attractiveness, condition, campus distance, and freshness.", accent=BLUE)
    add_card(slide, Inches(4.72), Inches(1.45), Inches(3.9), Inches(4.6), "Smart Price Suggestion",
             "Compares the draft listing with similar items, removes outliers, and generates a suggested price range with a median-based midpoint.", accent=GOLD)
    add_card(slide, Inches(8.74), Inches(1.45), Inches(3.9), Inches(4.6), "Trust and Safety",
             "Computes trust score from ratings, completed meetups, response time, email verification, and cancellation rate. Admin tools handle reports and moderation.", accent=BLUE)

    # Slide 7
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "System Architecture", "Client-server flow and technology stack", 7)
    add_arch_box(slide, Inches(0.95), Inches(1.7), Inches(2.25), Inches(0.8), "User Browser", WHITE)
    add_arch_box(slide, Inches(3.6), Inches(1.7), Inches(2.45), Inches(0.8), "React Frontend", WHITE)
    add_arch_box(slide, Inches(6.45), Inches(1.7), Inches(2.5), Inches(0.8), "Express REST API", WHITE)
    add_arch_box(slide, Inches(9.4), Inches(0.95), Inches(2.55), Inches(0.8), "MongoDB", WHITE)
    add_arch_box(slide, Inches(9.4), Inches(2.15), Inches(2.55), Inches(0.8), "Cloudinary / Local Uploads", WHITE)
    add_arch_box(slide, Inches(9.4), Inches(3.35), Inches(2.55), Inches(0.8), "Email Service", WHITE)
    add_arrow(slide, Inches(3.18), Inches(2.1), Inches(3.58), Inches(2.1))
    add_arrow(slide, Inches(6.03), Inches(2.1), Inches(6.43), Inches(2.1))
    add_arrow(slide, Inches(8.95), Inches(2.1), Inches(9.38), Inches(1.35))
    add_arrow(slide, Inches(8.95), Inches(2.1), Inches(9.38), Inches(2.55))
    add_arrow(slide, Inches(8.95), Inches(2.1), Inches(9.38), Inches(3.75))
    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.75),
        Inches(11.5),
        Inches(1.6),
        [
            "Frontend: React, React Router, Axios, Bootstrap, Leaflet",
            "Backend: Node.js, Express.js, JWT authentication, service-based business logic",
            "Database: MongoDB with Mongoose models for User, Listing, Message, Meetup, Review, Report, and Notification",
        ],
        font_size=17,
    )

    # Slide 8
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "Main Database Collections", "Core data entities used by the system", 8)
    table = add_table(slide, Inches(0.75), Inches(1.4), Inches(11.85), Inches(4.8), 8, 2)
    headers = [("Collection", "Purpose")]
    rows = headers + [
        ("User", "Stores account data, role, rating, tracked course codes, and preferences."),
        ("Listing", "Stores item information, seller reference, status, and meetup location."),
        ("Message", "Stores conversations and message history between participants."),
        ("Meetup", "Stores meetup proposal, confirmation state, and completion status."),
        ("Review", "Stores post-transaction rating and review content."),
        ("Report", "Stores user-submitted reports for listings or accounts."),
        ("Notification", "Stores alerts for messages, reservations, meetups, reviews, and wishlist matches."),
    ]
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r_idx == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(13 if r_idx == 0 else 12)
                    run.font.bold = r_idx == 0
                    run.font.color.rgb = WHITE if r_idx == 0 else DARK
    add_textbox(slide, Inches(0.9), Inches(6.4), Inches(11.0), Inches(0.45),
                "Relationships are centered around the marketplace flow: users create listings, listings lead to messages and meetups, and completed meetups enable reviews.", size=14, color=SLATE)

    # Slide 9
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "User Interface Demo I", "Homepage and search experience", 9)
    add_picture_contain(slide, screenshots["home"], Inches(0.6), Inches(1.3), Inches(6.0), Inches(4.65))
    add_caption(slide, Inches(0.6), Inches(6.0), Inches(6.0), "Homepage with recommendation board and latest listings")
    add_picture_contain(slide, screenshots["search"], Inches(6.8), Inches(1.3), Inches(5.95), Inches(4.65))
    add_caption(slide, Inches(6.8), Inches(6.0), Inches(5.95), "Search page with filters and matching results")

    # Slide 10
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "User Interface Demo II", "Listing detail and smart pricing workflow", 10)
    add_picture_contain(slide, screenshots["detail"], Inches(0.6), Inches(1.3), Inches(5.95), Inches(4.65))
    add_caption(slide, Inches(0.6), Inches(6.0), Inches(5.95), "Listing detail page with trust score and transaction actions")
    add_picture_contain(slide, screenshots["form"], Inches(6.8), Inches(1.3), Inches(5.95), Inches(4.65))
    add_caption(slide, Inches(6.8), Inches(6.0), Inches(5.95), "Create listing page with live smart price suggestion")

    # Slide 11
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "Communication and Administration", "Messaging, meetup planning, and moderation", 11)
    add_picture_contain(slide, screenshots["messages"], Inches(0.45), Inches(1.45), Inches(4.1), Inches(4.2))
    add_caption(slide, Inches(0.45), Inches(5.8), Inches(4.1), "Messages page")
    add_picture_contain(slide, screenshots["meetup"], Inches(4.63), Inches(1.45), Inches(4.1), Inches(4.2))
    add_caption(slide, Inches(4.63), Inches(5.8), Inches(4.1), "Smart meetup planner")
    add_picture_contain(slide, screenshots["admin"], Inches(8.8), Inches(1.45), Inches(4.1), Inches(4.2))
    add_caption(slide, Inches(8.8), Inches(5.8), Inches(4.1), "Admin dashboard")

    # Slide 12
    slide = new_slide(prs, blank)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)
    add_header(slide, "Testing, Deployment, and Conclusion", "Final validation and future direction", 12)
    add_card(slide, Inches(0.7), Inches(1.4), Inches(3.9), Inches(2.35), "Testing",
             "Backend automated tests: 11 passed, 0 failed.\nFrontend production build completed successfully.\nKey tested areas: reservation rules, search normalization, pricing, recommendations, notifications.", accent=BLUE)
    add_card(slide, Inches(4.75), Inches(1.4), Inches(3.9), Inches(2.35), "Deployment",
             "GitHub repository: github.com/Kise75/Group6-EduShare\nDemo method: run locally using the README guide.\nFrontend local URL: 127.0.0.1:5173\nBackend local URL: 127.0.0.1:5000", accent=GOLD)
    add_card(slide, Inches(8.8), Inches(1.4), Inches(3.85), Inches(2.35), "Conclusion",
             "EduShare is more than a simple CRUD system. It combines marketplace logic, safer meetup coordination, trust features, recommendation support, and admin moderation into one student-focused product.", accent=BLUE)
    add_bullets(
        slide,
        Inches(0.9),
        Inches(4.25),
        Inches(11.6),
        Inches(1.8),
        [
            "Future work: real-time messaging, more advanced recommendation personalization, image compression, and deployment automation.",
            "The project demonstrates both frontend and backend development skills with practical campus-oriented problem solving.",
        ],
        font_size=19,
    )
    add_textbox(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.35),
                "Thank you", size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    saved_main = save_with_fallback(prs, OUTPUT_PPTX)
    saved_english = save_with_fallback(prs, OUTPUT_PPTX_EN)
    print(saved_main)
    print(saved_english)


if __name__ == "__main__":
    build()
